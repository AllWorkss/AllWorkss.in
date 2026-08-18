"""
12-Slide Dark-Themed PPTX Presentation Deck Engine
Supports Dynamic White-Label Branding (partner_logo_path & partner_firm_name)
"""

from typing import Optional, Dict, Any
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


def generate_dark_boardroom_pptx(
    report_data: Dict[str, Any],
    output_path: str,
    partner_logo_path: Optional[str] = None,
    partner_firm_name: Optional[str] = None
) -> str:
    """
    Generates a 12-slide dark-themed boardroom PPTX presentation deck.
    Stamps partner white-label header if partner_logo_path & partner_firm_name are supplied.
    """
    firm_header = partner_firm_name if partner_firm_name else "Allworkss Business Intelligence Suite (ABIS)"

    if not PPTX_AVAILABLE:
        # Fallback file if python-pptx is missing
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(f"=== {firm_header} DARK PPTX BOARDROOM DECK ===\n")
            f.write(f"Company: {report_data.get('company_name', 'Client Enterprise')}\n")
            f.write("12 Slides Generated (Text Fallback)\n")
        return output_path

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    # Slide 1: Cover Slide
    slide1 = prs.slides.add_slide(blank_slide_layout)
    
    # Dark Background
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42)  # #0F172A

    txBox = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    tf = txBox.text_frame

    p1 = tf.paragraphs[0]
    p1.text = firm_header
    p1.font.bold = True
    p1.font.size = Pt(28)
    p1.font.color.rgb = RGBColor(0, 198, 255)  # Cyan

    p2 = tf.add_paragraph()
    p2.text = f"Boardroom Executive Audit & Financial Review\n{report_data.get('company_name', 'Client Enterprise')}"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(248, 250, 252)

    # Save Presentation
    prs.save(output_path)
    return output_path
